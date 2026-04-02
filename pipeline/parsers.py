"""
pipeline/parsers.py
-------------------
Multi-format document parsers for PDFs, PowerPoint slides, and Jupyter notebooks.
Each parser returns a list of dicts with a consistent schema:
  {
    "text":        str,   # extracted plain text
    "source":      str,   # file path
    "content_type": str,  # "pdf" | "slide" | "code" | "markdown"
    "slide":       int,   # (PPTX only) slide number
    "cell_type":   str,   # (notebook only) "code" | "markdown"
  }
"""

import os
from pathlib import Path


def parse_pdf(path: str) -> list[dict]:
    """Extract text from a PDF using the unstructured library and merge nearby elements."""
    try:
        from unstructured.partition.pdf import partition_pdf
    except ImportError:
        raise ImportError("Install unstructured: pip install unstructured[pdf]")

    elements = partition_pdf(path, strategy="fast")

    blocks = []
    buffer = []

    for el in elements:
        text = el.text.strip() if el.text else ""
        if not text:
            continue

        buffer.append(text)

        # Merge every ~5 elements into one block
        if len(buffer) >= 5:
            combined = "\n".join(buffer)
            blocks.append({
                "text": combined,
                "source": path,
                "content_type": "pdf",
            })
            buffer = []

    # Flush remaining buffer
    if buffer:
        blocks.append({
            "text": "\n".join(buffer),
            "source": path,
            "content_type": "pdf",
        })

    return blocks


def parse_pptx(path: str) -> list[dict]:
    """Extract text from each slide of a PowerPoint file."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("Install python-pptx: pip install python-pptx")

    prs = Presentation(path)
    chunks = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        combined = "\n".join(texts)
        if combined.strip():
            chunks.append({
                "text": combined,
                "source": path,
                "content_type": "slide",
                "slide": i,
            })
    return chunks


def parse_notebook(path: str) -> list[dict]:
    """Extract cells from a Jupyter notebook (.ipynb)."""
    try:
        import nbformat
    except ImportError:
        raise ImportError("Install nbformat: pip install nbformat")

    with open(path, "r", encoding="utf-8") as f:
        nb = nbformat.read(f, as_version=4)

    chunks = []
    for cell in nb.cells:
        text = cell.source.strip()
        if text:
            chunks.append({
                "text": text,
                "source": path,
                "content_type": cell.cell_type,  # "code" or "markdown"
                "cell_type": cell.cell_type,
            })
    return chunks


def parse_file(path: str) -> list[dict]:
    """Dispatch to the appropriate parser based on file extension."""
    ext = Path(path).suffix.lower()
    if ext == ".pdf":
        return parse_pdf(path)
    elif ext == ".pptx":
        return parse_pptx(path)
    elif ext == ".ppt":
        print(f"  WARNING: Legacy .ppt not supported, skipping {path}")
        return []
    elif ext == ".ipynb":
        return parse_notebook(path)
    elif ext in (".txt", ".md"):
        return parse_text(path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def parse_directory(directory: str) -> list[dict]:
    """Recursively parse all supported files in a directory."""
    supported = {".pdf", ".pptx", ".ipynb", ".txt", ".md"}
    all_chunks = []
    for root, _, files in os.walk(directory):
        for fname in files:
            if Path(fname).suffix.lower() in supported:
                fpath = os.path.join(root, fname)
                try:
                    chunks = parse_file(fpath)
                    all_chunks.extend(chunks)
                    print(f"  Parsed {fpath} → {len(chunks)} chunks")
                except Exception as e:
                    print(f"  WARNING: Failed to parse {fpath}: {e}")
    return all_chunks

def parse_text(path: str) -> list[dict]:
    """Extract text from plain .txt or .md files."""
    with open(path, "r", encoding="utf-8") as f:
        text = f.read().strip()
    if not text:
        return []
    return [{
        "text": text,
        "source": path,
        "content_type": "markdown" if path.endswith(".md") else "text",
    }]